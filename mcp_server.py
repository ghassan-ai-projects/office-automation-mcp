#!/usr/bin/env python3
"""
Office Automation MCP Server

An MCP (Model Context Protocol) server that wraps python-pptx into a type-safe,
persistent API for generating and editing PowerPoint presentations.

Run with:
    uv run --with python-pptx --with fastmcp /path/to/mcp_server.py
"""

from __future__ import annotations

import json
import os
import subprocess
import tempfile
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from mcp.server.fastmcp import FastMCP

from brand import Nature, TEMPLATES, TEMPLATE_NAMES
from slides.title import add_title_slide
from slides.content import add_content_slide
from slides.cards import add_cards_slide
from slides.process import add_process_slide
from slides.quotes import add_quotes_slide
from slides.funscale import add_funscale_slide
from slides.cta import add_cta_slide
from slides.thanks import add_thanks_slide
from utils import add_bg, SLIDE_W, SLIDE_H, _hex_to_rgb

# ── Server setup ──────────────────────────────────────────────────────────

mcp = FastMCP("Office Automation Engine")

PRESENTATION_DIMENSIONS = (SLIDE_W, SLIDE_H)


# ── Helper functions ──────────────────────────────────────────────────────

def _open_pptx(path):
    """Load or create a presentation."""
    if os.path.exists(path):
        return Presentation(path)
    return Presentation()


def _ensure_slide_dimensions(prs):
    """Ensure 16:9 widescreen dimensions."""
    target_w = int(SLIDE_W)
    target_h = int(SLIDE_H)
    if prs.slide_width != target_w or prs.slide_height != target_h:
        prs.slide_width = target_w
        prs.slide_height = target_h


def _save_and_open(prs, path):
    """Save presentation and optionally open in PowerPoint."""
    prs.save(path)
    try:
        subprocess.run(
            ["open", "-a", "Microsoft PowerPoint", path],
            capture_output=True, timeout=5
        )
    except Exception:
        pass  # PowerPoint not available — that's fine


def _get_slide_layout(prs, index=6):
    """Get a blank slide layout (index 6 = blank in default template)."""
    layouts = prs.slide_layouts
    if index < len(layouts):
        return layouts[index]
    # Fallback to last available layout
    return layouts[-1]


def _extract_shape_info(shape):
    """Extract readable info from a shape for JSON serialization."""
    info = {
        "type": shape.shape_type.name if hasattr(shape, "shape_type") else "unknown",
        "left": round(shape.left / 914400, 3) if shape.left is not None else 0,
        "top": round(shape.top / 914400, 3) if shape.top is not None else 0,
        "width": round(shape.width / 914400, 3) if shape.width is not None else 0,
        "height": round(shape.height / 914400, 3) if shape.height is not None else 0,
        "text": None,
        "font_size": None,
        "font_color": None,
        "fill_color": None,
    }
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            texts = [r.text for r in para.runs]
            if texts:
                info["text"] = "".join(texts)
                run = para.runs[0]
                if run.font.size:
                    info["font_size"] = round(run.font.size.pt, 1)
                if run.font.color and run.font.color.rgb:
                    info["font_color"] = str(run.font.color.rgb)
                break
    if hasattr(shape, "fill") and shape.fill.type is not None:
        try:
            info["fill_color"] = str(shape.fill.fore_color.rgb)
        except Exception:
            pass
    return info


def _read_slide_structure(path):
    """Read complete structure of a presentation."""
    prs = Presentation(path)
    slides_info = []
    for idx, slide in enumerate(prs.slides):
        bg_color = None
        try:
            fill = slide.background.fill
            if fill.type is not None:
                bg_color = str(fill.fore_color.rgb)
        except Exception:
            pass
        shapes = [_extract_shape_info(sh) for sh in slide.shapes]
        slides_info.append({
            "index": idx,
            "shape_count": len(slide.shapes),
            "background_color": bg_color,
            "shapes": shapes,
        })
    return {
        "slide_count": len(prs.slides),
        "dimensions": {
            "width": round(prs.slide_width / 914400, 3),
            "height": round(prs.slide_height / 914400, 3),
        },
        "slides": slides_info,
        "file_size_bytes": os.path.getsize(path) if os.path.exists(path) else 0,
    }


def _apply_brand_to_slide(slide, template_cls):
    """Apply brand colors to all shapes on a slide."""
    add_bg(slide, template_cls.BG)
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Keep existing accent colors, but text goes to deep forest
                        if run.font.color and run.font.color.rgb:
                            current = str(run.font.color.rgb)
                            # Don't overwrite white or accent colors
                            if current in ("FFFFFF", "1D6F42", "0E7490", "155534"):
                                continue
                        run.font.color.rgb = template_cls.TEXT_P
        except Exception:
            pass


# ── MCP Tools ─────────────────────────────────────────────────────────────

@mcp.tool()
def create_deck(path: str, template: str = "nature") -> str:
    """Create an empty branded presentation with 16:9 widescreen dimensions.

    Args:
        path: Absolute output path for the .pptx file
        template: Template name — "nature" (default), "dark", or "light-professional"

    Returns:
        JSON with status, path, slide_count, and slides_created
    """
    prs = Presentation()
    _ensure_slide_dimensions(prs)

    # Ensure output directory exists
    out_dir = os.path.dirname(os.path.abspath(path))
    os.makedirs(out_dir, exist_ok=True)

    # Remove default blank slide if present
    xml_slides = prs.slides._sldIdLst
    blank_slides = list(xml_slides)
    for elem in blank_slides:
        xml_slides.remove(elem)

    # Apply brand background to any remaining slides
    template_cls = TEMPLATES.get(template, Nature)
    for slide in prs.slides:
        add_bg(slide, template_cls.BG)

    _save_and_open(prs, path)

    result = {
        "status": "ok",
        "path": os.path.abspath(path),
        "slide_count": len(prs.slides),
        "slides_created": len(prs.slides),
    }
    return json.dumps(result)


@mcp.tool()
def add_slide(path: str, slide_type: str, data: str, index: int = -1) -> str:
    """Add a typed slide to an existing presentation.

    Args:
        path: Absolute path to existing .pptx
        slide_type: One of: title, content_bullets, cards, process_flow, quote, fun_scale, cta, thanks
        data: JSON string of slide-specific content
        index: Position to insert (-1 = end)

    Returns:
        JSON with status, slide_count, slide_index
    """
    data_dict = json.loads(data) if isinstance(data, str) else data
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    layout = _get_slide_layout(prs)
    blank_slide = prs.slides.add_slide(layout)

    # Build the slide using the appropriate builder
    builders = {
        "title": add_title_slide,
        "content_bullets": add_content_slide,
        "cards": add_cards_slide,
        "process_flow": add_process_slide,
        "quote": add_quotes_slide,
        "fun_scale": add_funscale_slide,
        "cta": add_cta_slide,
        "thanks": add_thanks_slide,
    }

    builder = builders.get(slide_type)
    if builder is None:
        return json.dumps({
            "status": "error",
            "error": f"Unknown slide_type: {slide_type}. Valid: {list(builders.keys())}"
        })

    builder(blank_slide, data_dict)

    # Reorder if index is specified
    if 0 <= index < len(prs.slides) - 1:
        # Move the new slide to the desired position
        xml_slides = prs.slides._sldIdLst
        last_elem = xml_slides[-1]
        xml_slides.remove(last_elem)
        if index >= len(xml_slides):
            xml_slides.append(last_elem)
        else:
            ref_elem = xml_slides[index]
            xml_slides.insert(xml_slides.index(ref_elem), last_elem)

    _save_and_open(prs, path)

    return json.dumps({
        "status": "ok",
        "slide_count": len(prs.slides),
        "slide_index": index if index >= 0 else len(prs.slides) - 1,
    })


@mcp.tool()
def apply_brand(path: str, template: str = "nature", slide_index: int = -1) -> str:
    """Apply brand colors to a slide or entire deck.

    Args:
        path: Absolute path to .pptx
        template: Template name — "nature" (default) or "light-professional"
        slide_index: Index of slide to brand (-1 = all slides)

    Returns:
        JSON with status and slides_updated count
    """
    template_cls = TEMPLATES.get(template, Nature)
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    slides_to_brand = prs.slides if slide_index < 0 else [prs.slides[slide_index]]
    count = 0
    for slide in slides_to_brand:
        _apply_brand_to_slide(slide, template_cls)
        count += 1

    prs.save(path)

    return json.dumps({
        "status": "ok",
        "slides_updated": count,
    })


@mcp.tool()
def read_slide_structure(path: str) -> str:
    """Read the current structure of a presentation.

    Args:
        path: Absolute path to .pptx

    Returns:
        JSON with complete slide structure
    """
    try:
        result = _read_slide_structure(path)
        return json.dumps(result)
    except Exception as e:
        return json.dumps({"status": "error", "error": str(e)})


@mcp.tool()
def duplicate_slide(path: str, source_index: int, target_index: int = -1) -> str:
    """Duplicate a slide within a presentation.

    Args:
        path: Absolute path to .pptx
        source_index: Index of slide to duplicate
        target_index: Position for the duplicate (-1 = end)

    Returns:
        JSON with status and slide_count
    """
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    if source_index < 0 or source_index >= len(prs.slides):
        return json.dumps({
            "status": "error",
            "error": f"source_index {source_index} out of range (0-{len(prs.slides) - 1})"
        })

    # Duplicate by deep-copying the slide XML element
    source_slide = prs.slides[source_index]
    slide_layout = source_slide.slide_layout

    # Add a new blank slide to get the required XML structure
    new_slide = prs.slides.add_slide(slide_layout)

    import copy as pycopy

    # Find the cSld element (p:cSld in presentationml namespace)
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    source_cSld = source_slide._element.find(f"{{{NS_P}}}cSld")
    new_cSld = new_slide._element.find(f"{{{NS_P}}}cSld")

    if source_cSld is not None and new_cSld is not None:
        # Remove all child elements of new cSld
        for child in list(new_cSld):
            new_cSld.remove(child)
        # Deep-copy all children from source cSld (spTree, bg, etc.)
        for child in source_cSld:
            new_cSld.append(pycopy.deepcopy(child))

    # Reorder if needed
    if target_index >= 0 and target_index < len(prs.slides) - 1:
        xml_slides = prs.slides._sldIdLst
        last_elem = xml_slides[-1]
        xml_slides.remove(last_elem)
        ref_elem = xml_slides[target_index]
        xml_slides.insert(xml_slides.index(ref_elem), last_elem)

    prs.save(path)

    return json.dumps({
        "status": "ok",
        "slide_count": len(prs.slides),
    })


@mcp.tool()
def delete_slide(path: str, index: int) -> str:
    """Remove a slide from a presentation.

    Args:
        path: Absolute path to .pptx
        index: Index of slide to delete

    Returns:
        JSON with status and slide_count
    """
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    if index < 0 or index >= len(prs.slides):
        return json.dumps({
            "status": "error",
            "error": f"index {index} out of range (0-{len(prs.slides) - 1})"
        })

    xml_slides = prs.slides._sldIdLst
    sldId = xml_slides[index]
    xml_slides.remove(sldId)

    prs.save(path)

    return json.dumps({
        "status": "ok",
        "slide_count": len(prs.slides),
    })


@mcp.tool()
def reorder_slides(path: str, order: str) -> str:
    """Reorder slides by specifying the new order of indices.

    Args:
        path: Absolute path to .pptx
        order: JSON array of integers — new order by index, e.g. [3, 1, 2, 4]

    Returns:
        JSON with status and slide_count
    """
    order_list = json.loads(order) if isinstance(order, str) else order
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    n = len(prs.slides)
    if len(order_list) != n:
        return json.dumps({
            "status": "error",
            "error": f"order length {len(order_list)} != slide count {n}"
        })

    if not all(0 <= i < n for i in order_list):
        return json.dumps({
            "status": "error",
            "error": f"all indices must be in range 0-{n - 1}"
        })

    if sorted(order_list) != list(range(n)):
        return json.dumps({
            "status": "error",
            "error": "order must be a permutation of 0..n-1"
        })

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    reordered = [slides[i] for i in order_list]
    for elem in slides:
        xml_slides.remove(elem)
    for elem in reordered:
        xml_slides.append(elem)

    prs.save(path)

    return json.dumps({
        "status": "ok",
        "slide_count": len(prs.slides),
    })


@mcp.tool()
def set_slide_background(path: str, slide_index: int, color: str) -> str:
    """Set a specific slide's background color.

    Args:
        path: Absolute path to .pptx
        slide_index: Index of the slide
        color: Hex color string, e.g. "#F5F7F5"

    Returns:
        JSON with status and slide_index
    """
    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    if slide_index < 0 or slide_index >= len(prs.slides):
        return json.dumps({
            "status": "error",
            "error": f"slide_index {slide_index} out of range (0-{len(prs.slides) - 1})"
        })

    slide = prs.slides[slide_index]
    rgb = _hex_to_rgb(color)
    add_bg(slide, rgb)
    prs.save(path)

    return json.dumps({
        "status": "ok",
        "slide_index": slide_index,
    })


@mcp.tool()
def add_image(path: str, slide_index: int, image_path: str,
              left: float = 0.0, top: float = 0.0,
              width: float = 4.0, height: float = 3.0) -> str:
    """Add an image to a slide.

    Args:
        path: Absolute path to .pptx
        slide_index: Index of the target slide
        image_path: Absolute path to the image file
        left: Left position in inches (default 0)
        top: Top position in inches (default 0)
        width: Width in inches (default 4)
        height: Height in inches (default 3)

    Returns:
        JSON with status
    """
    if not os.path.exists(image_path):
        return json.dumps({"status": "error", "error": f"Image not found: {image_path}"})

    prs = Presentation(path)
    _ensure_slide_dimensions(prs)

    if slide_index < 0 or slide_index >= len(prs.slides):
        return json.dumps({
            "status": "error",
            "error": f"slide_index {slide_index} out of range (0-{len(prs.slides) - 1})"
        })

    slide = prs.slides[slide_index]
    slide.shapes.add_picture(
        image_path,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    prs.save(path)

    return json.dumps({"status": "ok", "image_added": os.path.basename(image_path)})


@mcp.tool()
def export_as_pdf(path: str, output_path: str = "") -> str:
    """Export the deck as PDF.

    Since python-pptx cannot directly export to PDF, this uses macOS
    built-in printing (via osascript) to generate the PDF.

    Args:
        path: Absolute path to .pptx
        output_path: Desired PDF output path (default: same name with .pdf)

    Returns:
        JSON with status and pdf_path
    """
    if not os.path.exists(path):
        return json.dumps({"status": "error", "error": f"File not found: {path}"})

    abs_path = os.path.abspath(path)
    if not output_path:
        output_path = os.path.splitext(abs_path)[0] + ".pdf"

    output_abs = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_abs), exist_ok=True)

    try:
        # Use macOS scripting to open and export
        script = f'''
        tell application "Microsoft PowerPoint"
            set theDoc to open "{abs_path}"
            save theDoc in "{output_abs}" as "PDF"
            close theDoc
        end tell
        '''
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            # Fallback: use printing
            script2 = f'''
            tell application "Microsoft PowerPoint"
                set theDoc to open "{abs_path}"
                set pdfPath to "{output_abs}"
                -- Use save as PDF via Apple's virtual printer if needed
                close theDoc saving no
            end tell
            '''
            subprocess.run(
                ["osascript", "-e", script2],
                capture_output=True, timeout=30
            )
            # Try cups-pdf or manual approach
            subprocess.run(
                ["/usr/sbin/cupsfilter", abs_path, "-o", f"destination={output_abs}"],
                capture_output=True, timeout=30
            )
    except Exception as e:
        return json.dumps({
            "status": "warning",
            "message": f"Could not export via PowerPoint: {e}",
            "pdf_path": output_abs if os.path.exists(output_abs) else "",
        })

    exists = os.path.exists(output_abs)
    return json.dumps({
        "status": "ok" if exists else "warning",
        "pdf_path": output_abs if exists else "",
    })


# ── MCP Resources ─────────────────────────────────────────────────────────

@mcp.resource("slides://{path}")
def slides_structure(path: str) -> str:
    """Full slide structure as JSON — same as read_slide_structure."""
    try:
        return json.dumps(_read_slide_structure(path))
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.resource("slides://{path}/{index}")
def single_slide(path: str, index: str) -> str:
    """Single slide structure as JSON."""
    try:
        idx = int(index)
        full = _read_slide_structure(path)
        slides = full.pop("slides", [])
        if idx < 0 or idx >= len(slides):
            return json.dumps({"error": f"Slide index {idx} out of range"})
        result = {**full, "slide": slides[idx]}
        return json.dumps(result)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.resource("brand://current")
def brand_current() -> str:
    """Current brand configuration as JSON."""
    return json.dumps({
        "name": Nature.NAME,
        "colors": Nature.as_dict(),
        "fonts": {"primary": Nature.FONT},
        "templates_available": TEMPLATE_NAMES,
    })


# ── Entry point ───────────────────────────────────────────────────────────

def main():
    mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
