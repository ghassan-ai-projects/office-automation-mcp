#!/usr/bin/env python3
"""
Full 7-slide pitch deck example using the Office Automation MCP server.

This generates a complete client pitch deck on the local filesystem,
then uses the Python-pptx library directly for maximum control.
"""

import json
import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Add parent dir to path so we can import the modules directly
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from brand import Nature
from slides.title import add_title_slide
from slides.content import add_content_slide
from slides.cards import add_cards_slide
from slides.process import add_process_slide
from slides.quotes import add_quotes_slide
from slides.cta import add_cta_slide
from slides.thanks import add_thanks_slide
from utils import add_bg, SLIDE_W, SLIDE_H
from mcp_server import _ensure_slide_dimensions, _get_slide_layout, _save_and_open


def build_pitch_deck(output_path: str):
    """Generate a full 7-slide pitch deck."""
    prs = Presentation()
    _ensure_slide_dimensions(prs)

    blank_layout = _get_slide_layout(prs)

    # ── Slide 1: Title ─────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    add_title_slide(slide1, {
        "headline": "Velora Analytics",
        "subheadline": "Enterprise AI for Supply Chain Intelligence",
        "tagline": "Reduce waste · Predict demand · Optimize logistics",
        "attribution": "Confidential · Q3 2026",
        "url": "velora-analytics.io",
    })

    # ── Slide 2: Content / Problem ─────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    add_content_slide(slide2, {
        "headline": "The Problem",
        "subtitle": "Modern supply chains are drowning in data but starving for insight",
        "bullets": [
            "78% of supply chain leaders cite poor data visibility as their top challenge",
            "Legacy forecasting misses 35% of demand shifts",
            "Average enterprise loses $48M/year to logistics inefficiencies",
            "Manual analysis can't keep pace with real-time market changes",
        ],
        "footer": "Source: Gartner Supply Chain Survey 2025",
    })

    # ── Slide 3: Cards / Solution ──────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    add_cards_slide(slide3, {
        "headline": "The Velora Solution",
        "subtitle": "Four capabilities, one unified platform",
        "cards": [
            {
                "title": "Predictive Demand",
                "description": "ML models trained on 12+ years of supply chain data, achieving 94% forecast accuracy across 40+ industries.",
                "tags": ["AI", "Forecasting"],
            },
            {
                "title": "Real-time Visibility",
                "description": "Live dashboard with multi-echelon inventory tracking. Alerts when disruptions are detected — before they happen.",
                "tags": ["Monitoring", "IoT"],
            },
            {
                "title": "Cost Optimization",
                "description": "Route and load optimization algorithms that reduce shipping costs by up to 23%. Carbon-aware routing included.",
                "tags": ["Optimization", "ROI"],
            },
            {
                "title": "Compliance Engine",
                "description": "Automated regulatory compliance checks across 15+ global trade frameworks. Audit-ready reporting in one click.",
                "tags": ["Compliance", "Automation"],
            },
        ],
    })

    # ── Slide 4: Process Flow ──────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    add_process_slide(slide4, {
        "headline": "How It Works",
        "subtitle": "From integration to insight in four simple steps",
        "steps": [
            {"number": 1, "label": "Connect", "description": "API-first integration with your existing ERP, WMS, and TMS systems. Setup takes hours, not months."},
            {"number": 2, "label": "Ingest", "description": "Real-time data ingestion from 200+ sources. Cleansing and normalization happens automatically."},
            {"number": 3, "label": "Analyze", "description": "Proprietary ML models analyze patterns, detect anomalies, and generate actionable recommendations."},
            {"number": 4, "label": "Act", "description": "Push recommendations to your teams via Slack, email, or direct API. Close the loop with performance tracking."},
        ],
    })

    # ── Slide 5: Quotes / Testimonials ──────────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    add_quotes_slide(slide5, {
        "headline": "Trusted by Industry Leaders",
        "subtitle": "What our partners say about Velora Analytics",
        "quotes": [
            {
                "text": "Velora transformed our supply chain from a cost center to a competitive advantage. We saw 31% fewer stockouts in the first quarter.",
                "attribution": "Sarah Chen, VP Supply Chain · OmniCorp",
            },
            {
                "text": "The predictive demand engine alone saved us $12M in the first year. The ROI was immediate and continues to compound.",
                "attribution": "Marcus Johansson, COO · Nordic Logistics",
            },
        ],
    })

    # ── Slide 6: CTA ───────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    add_cta_slide(slide6, {
        "headline": "Ready to Transform Your Supply Chain?",
        "subheadline": "Schedule a demo and see Velora Analytics in action with your data",
        "button_text": "Book a Demo",
        "url": "velora-analytics.io/demo",
        "contact": "Your Name · email@example.com · +1 (555) 000-0000",
    })

    # ── Slide 7: Thanks ─────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    add_thanks_slide(slide7, {
        "headline": "Thank You",
        "tagline": "We look forward to partnering with you",
        "url": "velora-analytics.io",
    })

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ Pitch deck saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    out_dir = os.path.expanduser("~/ai-documents/presentations")
    output = os.path.join(out_dir, "velora-pitch-deck-2026-05-30.pptx")
    build_pitch_deck(output)
