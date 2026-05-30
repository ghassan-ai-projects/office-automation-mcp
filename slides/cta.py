"""CTA / closing slide builder."""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from brand import Nature
from utils import add_bg, add_accent_bar, add_textbox, add_shape, SLIDE_W, SLIDE_H


def add_cta_slide(slide, data):
    """Build a CTA/closing action slide.

    Schema: {headline, subheadline, button_text, url, contact}
    """
    # Background
    add_bg(slide, Nature.BG)

    # Top accent bar
    add_accent_bar(slide, left=Inches(0), top=Inches(0), width=SLIDE_W, height=Inches(0.08),
                   color=Nature.ACCENT)

    # Large headline
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
                text=data.get("headline", ""), font_size=36, font_color=Nature.TEXT_P,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Subheadline
    subhead = data.get("subheadline", "")
    if subhead:
        add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.6),
                    text=subhead, font_size=18, font_color=Nature.TEXT_S,
                    alignment=PP_ALIGN.CENTER)

    # CTA Button shape
    button_text = data.get("button_text", "")
    if button_text:
        btn_w = Inches(3)
        btn_h = Inches(0.55)
        btn_x = (SLIDE_W - btn_w) / 2
        btn_y = Inches(3.6)
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, btn_x, btn_y,
                                      btn_w, btn_h)
        btn.fill.solid()
        btn.fill.fore_color.rgb = Nature.ACCENT
        btn.line.fill.background()
        btn.adjustments[0] = 0.15
        tf = btn.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = button_text
        run.font.size = Pt(16)
        run.font.color.rgb = Nature.WHITE
        run.font.bold = True
        run.font.name = Nature.FONT
        tf._txBody.bodyPr.set("anchor", "ctr")

    # URL
    url = data.get("url", "")
    if url:
        add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
                    text=url, font_size=11, font_color=Nature.ACCENT,
                    alignment=PP_ALIGN.CENTER)

    # Contact
    contact = data.get("contact", "")
    if contact:
        add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4),
                    text=contact, font_size=12, font_color=Nature.TEXT_M,
                    alignment=PP_ALIGN.CENTER)

    # Bottom accent bar
    add_accent_bar(slide, left=Inches(0), top=SLIDE_H - Inches(0.06), width=SLIDE_W,
                   height=Inches(0.06), color=Nature.ACCENT)
