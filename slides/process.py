"""4-step horizontal process flow slide builder."""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from brand import Nature
from utils import add_bg, add_accent_bar, add_textbox, SLIDE_W, SLIDE_H


def add_process_slide(slide, data):
    """Build a 4-step horizontal process flow slide.

    Schema: {headline, subtitle, steps: [{number, label, description}]}
    Steps are displayed horizontally with connecting arrows.
    """
    # Background
    add_bg(slide, Nature.BG)

    # Top accent bar
    add_accent_bar(slide, left=Inches(0), top=Inches(0), width=SLIDE_W, height=Inches(0.06),
                   color=Nature.ACCENT)

    # Headline
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.5),
                text=data.get("headline", ""), font_size=28, font_color=Nature.ACCENT,
                bold=True)

    # Subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.35),
                    text=subtitle, font_size=14, font_color=Nature.TEXT_S)

    steps = data.get("steps", [])

    if not steps:
        return

    # Step dimensions
    num_steps = min(len(steps), 4)
    gap = Inches(0.15)
    arrow_w = Inches(0.4)
    step_w = (SLIDE_W - Inches(1.6) - (num_steps - 1) * (gap + arrow_w)) / num_steps
    step_h = Inches(3.5)
    start_y = Inches(1.5)
    start_x = Inches(0.8)

    for i in range(num_steps):
        s = steps[i]
        x = start_x + i * (step_w + gap + arrow_w)

        # Step card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y,
                                       step_w, step_h)
        card.fill.solid()
        card.fill.fore_color.rgb = Nature.CARD_BG
        card.line.color.rgb = Nature.BORDER
        card.line.width = Pt(1)
        card.adjustments[0] = 0.05

        # Step number circle
        circle_size = Inches(0.5)
        circ_x = x + (step_w - circle_size) / 2
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, start_y + Inches(0.15),
                                         circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = Nature.ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(s.get("number", i + 1))
        run.font.size = Pt(18)
        run.font.color.rgb = Nature.WHITE
        run.font.bold = True
        run.font.name = Nature.FONT
        tf._txBody.bodyPr.set("anchor", "ctr")

        # Step label
        add_textbox(slide, x + Inches(0.1), start_y + Inches(0.75), step_w - Inches(0.2), Inches(0.5),
                    text=s.get("label", ""), font_size=15, font_color=Nature.ACCENT,
                    bold=True, alignment=PP_ALIGN.CENTER)

        # Step description
        add_textbox(slide, x + Inches(0.1), start_y + Inches(1.3), step_w - Inches(0.2), Inches(2),
                    text=s.get("description", ""), font_size=11, font_color=Nature.TEXT_S,
                    alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < num_steps - 1:
            arrow_x = x + step_w + gap
            arrow_y = start_y + step_h / 2 - Inches(0.06)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y,
                                            arrow_w, Inches(0.12))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = Nature.ACCENT2
            arrow.line.fill.background()

    # Bottom accent bar
    add_accent_bar(slide, left=Inches(0), top=SLIDE_H - Inches(0.06), width=SLIDE_W,
                   height=Inches(0.06), color=Nature.ACCENT)
