"""Shared helper functions for building PPTX slides."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn, nsmap
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from brand import Nature, load_config

# ── Slide dimensions (configurable via themes.json / PPTX_THEME) ─────────
_config = load_config()
SLIDE_W = Inches(_config.get("slide_width_inches", 13.333))
SLIDE_H = Inches(_config.get("slide_height_inches", 7.5))
DEFAULT_FONT = _config.get("default_font", "Calibri")
FOOTER_TEXT = _config.get("footer_text", "")
AUTO_SAVE = _config.get("auto_save", True)

# ── Margins ───────────────────────────────────────────────────────────────
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
MARGIN_T = Inches(0.6)
MARGIN_B = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_bg(slide, color=None):
    """Set solid background fill on a slide."""
    if color is None:
        color = Nature.BG
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=0, top=Inches(0), width=None, height=Inches(0.06), color=None):
    """Add a horizontal accent bar (thin rectangle) across the slide."""
    if color is None:
        color = Nature.ACCENT
    if width is None:
        width = SLIDE_W
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=None, anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    if font_color is None:
        font_color = Nature.TEXT_P
    if font_name is None:
        font_name = Nature.FONT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b",
    }.get(anchor, "t"))
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_card(slide, left, top, width, height, title="", description="",
             tags=None, title_size=14, desc_size=11, tag_size=9,
             card_color=None, border_color=None, title_color=None,
             desc_color=None):
    """Add a rounded-rectangle card with title, description, and optional tags."""
    if card_color is None:
        card_color = Nature.CARD_BG
    if border_color is None:
        border_color = Nature.BORDER
    if title_color is None:
        title_color = Nature.ACCENT
    if desc_color is None:
        desc_color = Nature.TEXT_S
    if tags is None:
        tags = []

    # Card background rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    # Adjust corner radius
    shape.adjustments[0] = 0.05

    # Title
    pad = Inches(0.15)
    title_box = add_textbox(slide, left + pad, top + pad, width - pad * 2, Inches(0.4),
                            text=title, font_size=title_size, font_color=title_color,
                            bold=True)

    # Description
    desc_top = top + pad + Inches(0.35)
    desc_box = add_textbox(slide, left + pad, desc_top, width - pad * 2, height - Inches(0.5),
                           text=description, font_size=desc_size, font_color=desc_color)

    # Tag pills
    if tags:
        tag_left = left + pad
        tag_top = top + height - Inches(0.35)
        for tag in tags:
            pill_w = Inches(len(tag) * 0.07 + 0.25)
            add_tag_pill(slide, tag_left, tag_top, pill_w, Inches(0.22), tag, tag_size)
            tag_left += pill_w + Inches(0.06)

    return shape


def add_tag_pill(slide, left, top, width, height, text="", font_size=8, color=None):
    """Add a small tag pill (rounded rect with text)."""
    if color is None:
        color = Nature.ACCENT
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.3

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = Nature.WHITE
    run.font.name = Nature.FONT
    run.font.bold = False

    # Center vertically
    tf._txBody.bodyPr.set("anchor", "ctr")
    return shape


def add_shape(slide, shape_type, left, top, width, height, fill_color=None,
              line_color=None, line_width=None):
    """Add a basic shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape
